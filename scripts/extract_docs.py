#!/usr/bin/env python3
"""Phase 2 helper — extract operator-provided documents to plain text.

Operator documents (pitch decks, one-pagers, case-study PDFs, pricing sheets, battlecards) are the
highest-trust source for the company context. This script turns them into plain text so the synthesis
prompt can read them and a research subagent can be handed an excerpt.

Stdlib only for TXT / MD / CSV. For PDF / XLSX / PPTX / DOCX it TRIES an optional library inside a
try/except; if the library is not installed it prints a clear note telling the host agent to read that
file directly with its own file tools, and continues. It never hard-crashes on a missing optional
dependency, and it never requires a pip install to run.

Output: one `<name>.txt` per input under `<workdir>/docs_text/`. The working dir is resolved the same
way every script in this skill resolves it: the AUDIT_WORKDIR env var or a --workdir argument
(see scripts/lib/audit.py).

CLI:
    python3 scripts/extract_docs.py <indir-or-files...> --out <workdir>/docs_text
    python3 scripts/extract_docs.py raw/docs --workdir /abs/run            # --out defaults to <wd>/docs_text
    AUDIT_WORKDIR=/abs/run python3 scripts/extract_docs.py raw/docs
"""
import os
import re
import csv
import sys
import glob

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from lib import audit as A

TEXT_EXTS = {".txt", ".md", ".markdown", ".text", ".rst", ".log"}
CSV_EXTS = {".csv", ".tsv"}
PDF_EXTS = {".pdf"}
XLSX_EXTS = {".xlsx", ".xlsm"}
PPTX_EXTS = {".pptx"}
DOCX_EXTS = {".docx"}
KNOWN_EXTS = TEXT_EXTS | CSV_EXTS | PDF_EXTS | XLSX_EXTS | PPTX_EXTS | DOCX_EXTS


# ---------------------------------------------------------------------------- stdlib extractors

def extract_text(path):
    """Plain text / markdown / rst: read as-is. Stdlib only."""
    with open(path, encoding="utf-8", errors="replace") as f:
        return f.read()


def extract_csv(path):
    """CSV / TSV -> readable rows. Stdlib only. Delimiter sniffed, falls back to comma/tab by ext."""
    ext = os.path.splitext(path)[1].lower()
    with open(path, newline="", encoding="utf-8-sig", errors="replace") as f:
        sample = f.read(4096)
        f.seek(0)
        delim = "\t" if ext == ".tsv" else ","
        try:
            delim = csv.Sniffer().sniff(sample, delimiters=",\t;|").delimiter
        except Exception:
            pass
        rows = list(csv.reader(f, delimiter=delim))
    lines = []
    for r in rows:
        cells = [(c or "").strip() for c in r]
        if any(cells):
            lines.append(" | ".join(cells))
    return "\n".join(lines)


# ---------------------------------------------------------------------------- optional extractors
# Each returns (text, None) on success or (None, note) when the optional library is missing.
# Never raises on a missing dependency.

def extract_pdf(path):
    try:
        import pdfplumber  # type: ignore
    except ImportError:
        pdfplumber = None
    if pdfplumber is not None:
        try:
            out = []
            with pdfplumber.open(path) as pdf:
                for page in pdf.pages:
                    out.append(page.extract_text() or "")
            return "\n\n".join(out), None
        except Exception as e:
            return None, _read_directly_note(path, f"pdfplumber failed: {e}")
    # second-choice library
    PdfReader = None
    try:
        from pypdf import PdfReader  # type: ignore
    except ImportError:
        try:
            from PyPDF2 import PdfReader  # type: ignore
        except ImportError:
            PdfReader = None
    if PdfReader is not None:
        try:
            reader = PdfReader(path)
            return "\n\n".join((pg.extract_text() or "") for pg in reader.pages), None
        except Exception as e:
            return None, _read_directly_note(path, f"pypdf failed: {e}")
    return None, _read_directly_note(
        path, "no PDF library installed (tried pdfplumber, pypdf, PyPDF2)")


def extract_xlsx(path):
    try:
        import openpyxl  # type: ignore
    except ImportError:
        return None, _read_directly_note(path, "openpyxl not installed")
    try:
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    except Exception as e:
        return None, _read_directly_note(path, f"openpyxl failed: {e}")
    out = []
    for ws in wb.worksheets:
        out.append(f"# sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = ["" if c is None else str(c).strip() for c in row]
            if any(cells):
                out.append(" | ".join(cells))
        out.append("")
    return "\n".join(out), None


def extract_pptx(path):
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        return None, _read_directly_note(path, "python-pptx not installed")
    try:
        prs = Presentation(path)
    except Exception as e:
        return None, _read_directly_note(path, f"python-pptx failed: {e}")
    out = []
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"# slide {i}")
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for para in shape.text_frame.paragraphs:
                    txt = "".join(run.text for run in para.runs).strip()
                    if txt:
                        out.append(txt)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        out.append(" | ".join(cells))
        out.append("")
    return "\n".join(out), None


def extract_docx(path):
    try:
        from docx import Document  # type: ignore
    except ImportError:
        return None, _read_directly_note(path, "python-docx not installed")
    try:
        doc = Document(path)
    except Exception as e:
        return None, _read_directly_note(path, f"python-docx failed: {e}")
    out = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
    for tbl in doc.tables:
        for row in tbl.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                out.append(" | ".join(cells))
    return "\n".join(out), None


def _read_directly_note(path, reason):
    return (f"[NEEDS HOST AGENT] Could not extract {os.path.basename(path)} "
            f"({reason}). The host agent should read this file directly with its own file-reading "
            f"tool and paste the text into the company-research synthesis input. Do not skip it: "
            f"operator documents are the highest-trust source.")


# ---------------------------------------------------------------------------- dispatch + io

def out_name(path, used_names):
    """Deterministic, collision-free output stem: <slug>.txt, suffixing -2, -3 on collision."""
    stem = A.slug(os.path.splitext(os.path.basename(path))[0])
    name = stem
    n = 2
    while name in used_names:
        name = f"{stem}-{n}"
        n += 1
    used_names.add(name)
    return name + ".txt"


def gather_inputs(args):
    """Expand directories to their files (recursively) and keep listed files. Skips hidden + the
    output dir's own .txt artifacts so re-runs are idempotent."""
    paths = []
    for a in args:
        a = os.path.abspath(os.path.expanduser(a))
        if os.path.isdir(a):
            for root, dirs, files in os.walk(a):
                dirs[:] = [d for d in dirs if not d.startswith(".")]
                for fn in sorted(files):
                    if fn.startswith("."):
                        continue
                    paths.append(os.path.join(root, fn))
        elif os.path.isfile(a):
            paths.append(a)
        else:
            for g in sorted(glob.glob(a)):
                if os.path.isfile(g):
                    paths.append(g)
    # de-dupe, preserve order
    seen, ordered = set(), []
    for p in paths:
        if p not in seen:
            seen.add(p)
            ordered.append(p)
    return ordered


def extract_one(path):
    """Returns (text, note). text is None when extraction needs the host agent (note explains)."""
    ext = os.path.splitext(path)[1].lower()
    try:
        if ext in TEXT_EXTS:
            return extract_text(path), None
        if ext in CSV_EXTS:
            return extract_csv(path), None
        if ext in PDF_EXTS:
            return extract_pdf(path)
        if ext in XLSX_EXTS:
            return extract_xlsx(path)
        if ext in PPTX_EXTS:
            return extract_pptx(path)
        if ext in DOCX_EXTS:
            return extract_docx(path)
    except Exception as e:  # never hard-crash on one bad file
        return None, _read_directly_note(path, f"unexpected error: {e}")
    # unknown extension: try to read it as text; if it is binary, hand it to the host agent
    try:
        txt = extract_text(path)
        if "\x00" in txt[:2048]:
            return None, _read_directly_note(path, f"unknown/binary extension '{ext}'")
        return txt, None
    except Exception as e:
        return None, _read_directly_note(path, f"unreadable ({e})")


def resolve_out_dir():
    """--out wins; else <workdir>/docs_text (workdir from --workdir or AUDIT_WORKDIR)."""
    argv = sys.argv
    for i, a in enumerate(argv):
        if a == "--out" and i + 1 < len(argv):
            return os.path.abspath(os.path.expanduser(argv[i + 1]))
        if a.startswith("--out="):
            return os.path.abspath(os.path.expanduser(a.split("=", 1)[1]))
    return os.path.join(A.workdir(), "docs_text")


def positional_args():
    """Everything that is not a flag or a flag's value."""
    argv = sys.argv[1:]
    out = []
    skip = False
    for i, a in enumerate(argv):
        if skip:
            skip = False
            continue
        if a in ("--out", "--workdir"):
            skip = True
            continue
        if a.startswith(("--out=", "--workdir=")):
            continue
        out.append(a)
    return out


def main():
    inputs = positional_args()
    if not inputs:
        print(__doc__)
        sys.exit("ERROR: pass at least one input directory or file.")
    out_dir = resolve_out_dir()
    os.makedirs(out_dir, exist_ok=True)

    files = gather_inputs(inputs)
    if not files:
        print(f"No input files found under: {', '.join(inputs)}")
        sys.exit(0)

    used = set()
    written, needs_agent = [], []
    for path in files:
        text, note = extract_one(path)
        target = os.path.join(out_dir, out_name(path, used))
        if text is not None and text.strip():
            header = f"# source: {os.path.basename(path)}\n# path: {path}\n\n"
            with open(target, "w", encoding="utf-8") as f:
                f.write(header + text.strip() + "\n")
            written.append((os.path.basename(path), target, len(text)))
        else:
            note = note or _read_directly_note(path, "no text extracted")
            needs_agent.append((path, note))

    print(f"\n================ DOC EXTRACTION — {len(files)} file(s) ================")
    print(f"  out dir: {out_dir}")
    if written:
        print(f"\n  extracted to text ({len(written)}):")
        for name, target, n in written:
            print(f"    {name:40s} -> {os.path.basename(target)}  ({n} chars)")
    if needs_agent:
        print(f"\n  NEEDS HOST AGENT — read these directly with your own file tools ({len(needs_agent)}):")
        for path, note in needs_agent:
            print(f"    {os.path.basename(path)}")
            print(f"      {note}")
        print("\n  (A missing optional library is expected. Install pdfplumber / openpyxl / "
              "python-pptx / python-docx to extract these here, OR just read them with the host "
              "agent's file reader. Either way the document is not dropped.)")
    if not written and not needs_agent:
        print("  nothing to extract.")
    print("=================================================================\n")
    return 0


if __name__ == "__main__":
    sys.exit(main())
